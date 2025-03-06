from django.shortcuts import render
from django.contrib.auth.decorators import login_required
from django.conf import settings
from django.http import JsonResponse
import json
import google.generativeai as genai
from django.views.decorators.csrf import csrf_exempt
import os
from PyPDF2 import PdfReader
import docx
from pptx import Presentation  # For PowerPoint support

# Configure Gemini API
genai.configure(api_key=settings.GEMINI_API_KEY)

# Store uploaded file text for each user
DOCUMENTS = {}

@login_required
def chatbot_page(request):
    """Renders the chatbot page."""
    return render(request, "chatbot/chatbot.html")

@csrf_exempt
@login_required
def chatbot_response(request):
    """Handles chatbot queries and responds with text-based or document-based answers."""
    if request.method == "POST":
        data = json.loads(request.body)
        user_query = data.get("query", "")
        request_type = data.get("type", "text")  # "text", "document", "summary"
        
        if not user_query and request_type != "summary":
            return JsonResponse({"error": "Query cannot be empty"}, status=400)

        try:
            model = genai.GenerativeModel("gemini-1.5-flash")
            document_text = DOCUMENTS.get(request.user.id, "")
            
            if request_type in ["document", "summary"]:
                if not document_text:
                    return JsonResponse({"error": "No document uploaded yet"}, status=400)
                
                query_with_context = (
                    f"Summarize the following text:\n\n{document_text}" if request_type == "summary" 
                    else f"Using the provided document, answer this question: {user_query}\n\nDocument:\n{document_text}"
                )
                
                response = model.generate_content(query_with_context)
                return JsonResponse({"response": response.text})
            
            # Default chatbot response
            response = model.generate_content(user_query)
            return JsonResponse({"response": response.text})

        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)

    return JsonResponse({"error": "Invalid request"}, status=400)

@csrf_exempt
@login_required
def upload_file(request):
    """Handles file uploads and extracts text from the document."""
    if request.method == "POST" and request.FILES.get("file"):
        uploaded_file = request.FILES["file"]
        file_ext = os.path.splitext(uploaded_file.name)[1].lower()
        text = ""
        
        try:
            if file_ext == ".pdf":
                pdf_reader = PdfReader(uploaded_file)
                text = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
            elif file_ext == ".docx":
                doc = docx.Document(uploaded_file)
                text = "\n".join([para.text for para in doc.paragraphs])
            elif file_ext == ".txt":
                text = uploaded_file.read().decode("utf-8")
            elif file_ext in [".ppt", ".pptx"]:
                ppt = Presentation(uploaded_file)
                slide_texts = []
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            slide_texts.append(shape.text_frame.text)
                text = "\n".join(slide_texts)
            else:
                return JsonResponse({"error": "Unsupported file format"}, status=400)

            # Store the extracted text for user-specific access
            DOCUMENTS[request.user.id] = text
            return JsonResponse({"message": "File uploaded successfully", "text_preview": text[:500] + "..."})
        
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)

    return JsonResponse({"error": "No file uploaded"}, status=400)
